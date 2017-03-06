#VERSION 1.     MAKES PEPPER SPEAK THE NOTES OF THE PRESENTATION
#               REMOTELY CHOSEN FROM THE PC WITH THE GUI

from naoqi import ALProxy
from pptx import Presentation
from Tkinter import Tk
from tkFileDialog import askopenfilename

robot_ip = "set robot ip" #!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!
#define presentation path with the GUI, form the computer
Tk().withdraw()
my_path = askopenfilename()

#create a presentation object from the library pptx and get its number of slides
my_pres = Presentation(my_path)
number_of_slides = len(my_pres.slides)

#crete the proxy for speech
tts = ALProxy("ALTextToSpeech", robot_ip, 9559)

#iterate over the number of slides, making pepper tell the notes
for slide in my_pres.slides:
        #iterate over the slide notes
        if slide.has_notes_slide: #chech if we have notes
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                #those 2 following extractions are due to how the XML is structured
                for paragraph in text_frame.paragraphs:
                        #print paragraph.text
                        tts.say(paragraph.text)