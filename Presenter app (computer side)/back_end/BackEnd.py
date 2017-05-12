import htmlPy
import os
import json
from naoqi import ALProxy
from pptx import Presentation
import time
import math

#linked class with front end
class PepperApp(htmlPy.Object):
	#GUI callable functions need to be in a class, inherited from htmlPy.Object and binded to the GUI
	def __init__(self, app):
		super(PepperApp, self).__init__()
		self.app = app
		self.my_path = ""
		self.my_ip = ""
		self.current_slide =""
		self.start_slide =""
		self.slides_path =""
		self.tts =""
		self.my_pres =""
		self.number_of_slides = ""
		self.gestures_dict = {":)":"happy", ":(":"sad", ":|":"unknown"}
		self.pointing_right = {"[point=1]":"uf", "[point=2]":"un", "[point=3]":"ln", "[point=4]":"lf"}
		self.pointing_left = {"[point=1]":"un", "[point=2]":"uf", "[point=3]":"lf", "[point=4]":"ln"}
		self.isRobotPointing = False
		self.leds = {":)":"green", ":(":"magenta", ":|":"yellow"}
	#functions are defined here, under the .Slot decoration
	
############################################################
############################################################
#
#				STARTING PRESENTATION (ON CLICKING START)
#				imports everything, shows first slide, calculates pointing angles
#
############################################################
############################################################
	@htmlPy.Slot(str)
	def start_presentation(self, json_data): #initializes everything
#retrieve pointing data and calculates turn angle
		def retrieve_pointing_data():
			self.width = float(form_data['screen-width'])
			self.height = float(form_data['screen-height'])
			self.x = float(form_data['x'])
			self.y = float(form_data['y'])
			self.z = float(form_data['z'])
			self.near = self.x + self.width/4
			self.far = self.x + 3*self.width/4
			self.upper = self.z + 3*self.height/4 -0.91
			self.lower = self.z + self.height/4 -0.91
			print(self.near)
			print(self.far)
			if form_data['side'] == '':
				self.side ='left'
				self.shoulder = 'LShoulderPitch'
				self.pointing_dict = self.pointing_left
				self.zNear = (math.atan(self.y/self.near)+math.pi/2)
				self.zFar = (math.atan(self.y/self.far)+math.pi/2)
				print(self.zNear)
				print(self.zFar)
			elif form_data['side'] =='on':
				self.side = 'right'	
				self.shoulder = 'RShoulderPitch'
				self.pointing_dict = self.pointing_right
				self.zNear = -(math.atan(self.y/self.near)+math.pi/2)
				self.zFar = -(math.atan(self.y/self.far)+math.pi/2)
			print(self.side)
			return
#calculates pointing angles vector  = {'un':[angleZ,angleShoulder],'uf':[angleZ,angleShoulder],'ln':[angleZ,angleShoulder], 'lf':[angleZ,angleShoulder]}
		def calculate_angles():
			self.upperNear = -math.atan(self.upper/(math.sqrt(math.pow(self.near,2)+math.pow(self.y,2))))
			self.upperFar = -math.atan(self.upper/(math.sqrt(math.pow(self.far,2)+math.pow(self.y,2))))
			self.lowerNear = -math.atan(self.lower/(math.sqrt(math.pow(self.near,2)+math.pow(self.y,2))))
			self.lowerFar = -math.atan(self.lower/(math.sqrt(math.pow(self.far,2)+math.pow(self.y,2))))

			self.angles = {'un':[self.zNear,self.upperNear],'uf':[self.zFar,self.upperFar],'ln':[self.zNear,self.lowerNear], 'lf':[self.zFar,self.lowerFar]}
			print(self.angles)
			return
#read the json data from the form (ip, file, start)
		form_data = json.loads(json_data)
		self.my_path = form_data['file']
		self.my_ip = form_data['ip']
		self.start_slide = form_data['start']
		self.point_check = form_data['enable_point'] 
		
		if self.point_check == '':
			self.point_enabled = False
#if gesture is enable then get the coordinates variables	
		elif self.point_check == 'on':
			self.point_enabled = True
			retrieve_pointing_data()
			calculate_angles()
#initialize the presentation and variables
		self.start_slide = int(self.start_slide)
		self.current_slide = self.start_slide
		self.slides_path = os.path.dirname(self.my_path) + '/slides/slide'
		self.my_pres = Presentation(self.my_path)
		self.number_of_slides = len(self.my_pres.slides)
		notes = []
#connect to the robot and show initial slide

		#COMENT THIS WHEN TESTING OUTISDE ROBOT
	########################################################################	
		self.tts = ALProxy("ALAnimatedSpeech", str(self.my_ip), 9559)
		self.motion = ALProxy("ALMotion", str(self.my_ip), 9559)
		self.posture = ALProxy("ALRobotPosture", str(self.my_ip), 9559)
		self.aw = ALProxy("ALBasicAwareness", str(self.my_ip), 9559)
		self.motion.moveInit()
		self.motion.setStiffnesses(self.shoulder, 1)
		self.posture.goToPosture("StandInit", 0.5)
		self.aw.setTrackingMode("Head")
		self.aw.setEngagementMode("Unengaged")
		# self.aw.resumeAwareness()
		self.aw.pauseAwareness()
	########################################################################
		slide_src = self.slides_path + str(self.start_slide) + '.jpg'
		self.app.evaluate_javascript("document.getElementById('presentation_image').style.display='block'")
		self.app.evaluate_javascript("document.getElementById('presentation_content').innerHTML = 'Log:<br>Starting presentation at: %s<br>IP: %s<br>Notes: '" %(self.my_path, self.my_ip))
		self.app.evaluate_javascript("document.getElementById('slide').src = '%s'" %(slide_src))
		self.app.evaluate_javascript("document.getElementById('presentation_image').style.display = 'block'")
		self.app.evaluate_javascript("scroll(0,0)")
		print('Showing slide ' + str(self.current_slide) +'. Source: '+ slide_src)

#the calculations of angles for the pointing function should be done here, define them as self variables to access from present_slide
		return

############################################################
############################################################
#
#				PRESENTING SLIDE (ON CLICK SLIDE)
#
############################################################
############################################################
	@htmlPy.Slot()
	def present_slide(self):
		self.aw.setTrackingMode("Head")
		self.aw.setEngagementMode("Unengaged")
#this will use the dictionary to check for gestures  FIX IT TO ONLY ADD ONE CLOSING TAG
		def check_gestures(text):
			for key in self.gestures_dict:
				if text.find(key) != -1:
					# text = text.replace(key, "^startTag(" + self.gestures_dict[key] + ")" ) + " ^stopTag(" + self.gestures_dict[key] + ")"
					#here we try to make the emotion at the end of the sentence		 
					#not sure if the double \\ is needed as escape character 
					text = '^call(ALLeds.fadeRGB("FaceLeds", "'+ self.leds[key] + '", 3)) '	+ text.replace(key, "" ) + ' \\pau=1000\\ '+ ' ^startTag(' + self.gestures_dict[key] + ')' + ' ^waitTag(' + self.gestures_dict[key] + ')'
			return text

		def check_point(text):
			for key in self.pointing_dict:
				if text.find(key) != -1:
					text = text.replace(key, "")					
					print('Pointing to ' + self.pointing_dict[key])
					#COMMENT THIS WHEN TESTING OUTISDE ROBOT
		####################################################################
					# self.aw.pauseAwareness()
		####################################################################
					#point to that position
					point(self.pointing_dict[key])
					
			return text

		def point(position):
			self.angles_vector = self.angles[position] #now we have the vector with the Z angle and the shoulder angle for a given position
			
			#COMMENT THIS WHEN TESTING OUTISDE ROBOT
		####################################################################	
			self.motion.setStiffnesses(self.shoulder, 1)
			self.motion.moveTo(0,0,self.angles_vector[0])
			self.motion.setAngles(self.shoulder, self.angles_vector[1],0.3)
		####################################################################
			self.isRobotPointing = True
			
			return

#with calls inside the text
################################################################################		
# 		def point(position, text):
# 			self.angles_vector = self.angles[position] #now we have the vector with the Z angle and the shoulder angle for a given position
# 			my_z_str = str(self.angles_vector[0])
# 			my_shoulder_str = str(self.angles_vector[1])
# #^pcall() = asynchronous, ^call() = ^synchronous
# 			animated_text = '^call(ALMotion.moveTo(0,0,'+ my_z_str+')) ^call(ALMotion.setAngles("' + self.shoulder + '",' + my_shoulder_str+',0.3))' + text
			
# 			self.isRobotPointing = True
# 			return animated_text

#ORiginal code,  the upper one is modified
################################################################################
# 		def check_point(text):
# 			for key in self.pointing_dict:
# 				if text.find(key) != -1:
# 					text = text.replace(key, "")					
# 					print('Pointing to ' + self.pointing_dict[key])
# 					#point to that position
# 					point(self.pointing_dict[key])
					
# 			return text
# #the function to point should be added here also
# 		def point(position):
# 			self.angles_vector = self.angles[position] #now we have the vector with the Z angle and the shoulder angle for a given position
			
		# 	#COMMENT THIS WHEN TESTING OUTISDE ROBOT
		# ####################################################################	
		# 	self.motion.setStiffnesses(self.shoulder, 1)
		# 	self.motion.moveTo(0,0,self.angles_vector[0])
		# 	self.motion.setAngles(self.shoulder, self.angles_vector[1],0.5)
		# ####################################################################
		# 	self.isRobotPointing = True
		# 	return
###########################################################################################


#the slide is showing, so when you click on it it will read the notes of the slide
#if it is not the last one it will show the next slide, if it is the last one will elapse some time and close the image view
		slide = self.my_pres.slides[self.current_slide-1]
		if slide.has_notes_slide:
			notes_slide = slide.notes_slide
			text_frame = notes_slide.notes_text_frame
			for paragraph in text_frame.paragraphs:
				if self.point_enabled:
					after_pointing_txt = check_point(paragraph.text)
				else:
					after_pointing_txt = paragraph.text
				modified_text = check_gestures(after_pointing_txt)
				self.app.evaluate_javascript("document.getElementById('presentation_content').innerHTML += '<br>  %s  -  %s '" %(paragraph.text, modified_text))
				print('Notes line of slide ' + str(self.current_slide) +': ' + paragraph.text)
				print('Modified notes line of slide ' + str(self.current_slide) +': ' + modified_text)

				#COMMENT THIS WHEN TESTING OUTISDE ROBOT
			####################################################################
				self.tts.say(str(modified_text))
			
				if self.isRobotPointing:
					self.motion.moveTo(0,0,-self.angles_vector[0], _async = True)
					self.isRobotPointing = False
					# self.aw.resumeAwareness()
			####################################################################
				time.sleep(0.5)
		self.current_slide +=1
		if self.current_slide <=self.number_of_slides:
			slide_src = self.slides_path + str(self.current_slide) + '.jpg'
			self.app.evaluate_javascript("document.getElementById('slide').src = '%s'" %(slide_src))
		else:
			time.sleep(2)
			#COMMENT THIS WHEN TESTING OUTISDE ROBOT
		####################################################################
			self.aw.resumeAwareness()
			self.aw.setTrackingMode("Head")
			self.aw.setEngagementMode("Unengaged")
		####################################################################

			self.app.evaluate_javascript("document.getElementById('presentation_image').style.display = 'none'")

		return

# #WORKING!

